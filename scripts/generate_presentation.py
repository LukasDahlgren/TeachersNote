import argparse
import json
import re
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO

import fitz  # pymupdf
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
IMAGE_PANEL_WIDTH = Inches(8.0)
NOTES_PANEL_LEFT = Inches(8.0)
NOTES_PANEL_WIDTH = Inches(5.33)
NOTES_PAD = Inches(0.18)
NOTES_BG = RGBColor(0xF5, 0xF7, 0xFA)
HEADER_COLOR = RGBColor(0x0F, 0x4C, 0x81)
BODY_COLOR = RGBColor(0x22, 0x29, 0x36)
BULLET_PREFIX_RE = re.compile(r"^\s*(?:[-*•]\s+|\d+[.)]\s+)")


def pdf_to_images(pdf_path: str, dpi: int = 150) -> list[tuple[bytes, float]]:
    """Render each PDF page to PNG bytes in parallel; return (png_bytes, aspect_ratio)."""
    doc = fitz.open(pdf_path)
    try:
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        page_count = len(doc)

        def render_page(idx: int) -> tuple[int, bytes, float]:
            page = doc[idx]
            aspect = page.rect.width / page.rect.height
            pix = page.get_pixmap(matrix=matrix)
            return idx, pix.tobytes("png"), aspect

        with ThreadPoolExecutor() as pool:
            results = list(pool.map(render_page, range(page_count)))
    finally:
        doc.close()
    return [(img, asp) for _, img, asp in sorted(results)]


def _bulletize_text(text: str) -> list[str]:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not normalized:
        return []

    lines = [line.strip() for line in normalized.split("\n") if line.strip()]
    if not lines:
        return []

    prefixed_items: list[str] = []
    has_prefixed = False
    for line in lines:
        match = BULLET_PREFIX_RE.match(line)
        if match:
            has_prefixed = True
            item = line[match.end():].strip()
            if item:
                prefixed_items.append(item)
            continue
        if has_prefixed and prefixed_items:
            prefixed_items[-1] = f"{prefixed_items[-1]} {line}".strip()

    if has_prefixed and prefixed_items:
        return prefixed_items

    if len(lines) > 1:
        return lines

    compact = " ".join(normalized.split())
    parts = [p.strip() for p in re.split(r"(?<=[.!?])\s+|;\s+", compact) if p.strip()]
    return parts or ([compact] if compact else [])


def build_speaker_notes(entry: dict) -> str:
    lines = []
    lines.append(f"SAMMANFATTNING: {entry.get('summary', '')}")
    lines.append("")
    la = entry.get("lecturer_additions", "")
    lecturer_bullets = _bulletize_text(la) if la else []
    if lecturer_bullets:
        lines.append("FÖRELÄSARENS TILLÄGG:")
        for item in lecturer_bullets:
            lines.append(f"  • {item}")
        lines.append("")
    takeaways = entry.get("key_takeaways", [])
    if takeaways:
        lines.append("KEY TAKEAWAYS:")
        for t in takeaways:
            lines.append(f"  • {t}")
    return "\n".join(lines)


def _place_image(slide, img_bytes: bytes, aspect: float) -> None:
    """Place the PDF page image in the left panel, scaled to fit, vertically centred."""
    max_w = IMAGE_PANEL_WIDTH
    max_h = SLIDE_HEIGHT
    if aspect > max_w / max_h:
        pic_w = max_w
        pic_h = int(max_w / aspect)
    else:
        pic_h = max_h
        pic_w = int(max_h * aspect)
    top = (max_h - pic_h) // 2
    slide.shapes.add_picture(BytesIO(img_bytes), 0, top, pic_w, pic_h)


def _add_notes_panel(slide, entry: dict) -> None:
    """Add a light right-panel textbox with structured notes sections."""
    txBox = slide.shapes.add_textbox(NOTES_PANEL_LEFT, Inches(0), NOTES_PANEL_WIDTH, SLIDE_HEIGHT)

    txBox.fill.solid()
    txBox.fill.fore_color.rgb = NOTES_BG

    tf = txBox.text_frame
    tf.word_wrap = True

    # Inner padding via body properties
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        pad_str = str(int(NOTES_PAD))
        bodyPr.set("lIns", pad_str)
        bodyPr.set("rIns", pad_str)
        bodyPr.set("tIns", pad_str)
        bodyPr.set("bIns", pad_str)

    def _add_header(label: str, first: bool = False) -> None:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(0) if first else Pt(4)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = HEADER_COLOR

    def _add_body(text: str) -> None:
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(8.5)
        run.font.color.rgb = BODY_COLOR

    def _add_bullet(text: str) -> None:
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"• {text}"
        run.font.size = Pt(8.5)
        run.font.color.rgb = BODY_COLOR

    # Section 1: SAMMANFATTNING
    summary = entry.get("summary", "").strip()
    _add_header("SAMMANFATTNING", first=True)
    if summary:
        _add_body(summary)

    # Section 2: FÖRELÄSARENS TILLÄGG
    la = entry.get("lecturer_additions", "")
    bullets = _bulletize_text(la) if la else []
    if bullets:
        _add_header("FÖRELÄSARENS TILLÄGG")
        for b in bullets:
            _add_bullet(b)

    # Section 3: KEY TAKEAWAYS
    takeaways = entry.get("key_takeaways", [])
    if takeaways:
        _add_header("KEY TAKEAWAYS")
        for t in takeaways:
            _add_bullet(t)


def generate(pdf_path: str, enhanced_path: str, output_path: str) -> None:
    with open(enhanced_path, encoding="utf-8") as f:
        enhanced = json.load(f)

    enhanced_by_slide = {e["slide"]: e for e in enhanced}

    print("Rendering PDF pages to images...")
    images = pdf_to_images(pdf_path)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for i, (img_bytes, aspect) in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank_layout)

        _place_image(slide, img_bytes, aspect)

        entry = enhanced_by_slide.get(i)
        if entry:
            _add_notes_panel(slide, entry)
            # Also keep speaker notes for compatibility
            notes_text = build_speaker_notes(entry)
            tf = slide.notes_slide.notes_text_frame
            tf.text = notes_text
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)

        print(f"  Slide {i}/{len(images)} done")

    prs.save(output_path)
    print(f"\nSaved {len(images)}-slide presentation → {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate enriched PPTX from PDF slides + enhanced.json"
    )
    parser.add_argument("--pdf", required=True, help="Original PDF lecture slides")
    parser.add_argument("--enhanced", required=True, help="Path to enhanced.json")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    parser.add_argument("--dpi", type=int, default=150, help="Render DPI (default 150)")
    args = parser.parse_args()
    generate(args.pdf, args.enhanced, args.output)
